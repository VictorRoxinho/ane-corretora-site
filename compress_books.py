import subprocess
import sys

# Auto-install dependencies
subprocess.check_call([sys.executable, "-m", "pip", "install", "pymupdf", "python-pptx", "python-docx", "Pillow", "-q"])

import os
import io
import fitz  # pymupdf
from PIL import Image
from pptx import Presentation
from pptx.util import Pt
from docx import Document
from pathlib import Path

BOOKS_DIR = Path(__file__).parent / "public" / "books"


def fmt_size(bytes_val):
    return f"{bytes_val / 1024 / 1024:.2f} MB"


def compress_pdf(src: Path, dst: Path):
    doc = fitz.open(src)
    # Recompress with reduced image quality
    for page in doc:
        for img in page.get_images(full=True):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            pil_img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
            buf = io.BytesIO()
            pil_img.save(buf, format="JPEG", quality=60, optimize=True)
            doc.update_stream(xref, buf.getvalue())
    doc.save(dst, garbage=4, deflate=True, clean=True)
    doc.close()


def compress_image_bytes(data: bytes, quality=60) -> bytes:
    try:
        img = Image.open(io.BytesIO(data))
        fmt = img.format or "JPEG"
        buf = io.BytesIO()
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")
        img.save(buf, format="JPEG", quality=quality, optimize=True)
        compressed = buf.getvalue()
        return compressed if len(compressed) < len(data) else data
    except Exception:
        return data


def compress_pptx(src: Path, dst: Path):
    prs = Presentation(src)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image = shape.image
                compressed = compress_image_bytes(image.blob)
                # Replace blob in-place via the part
                image._blob = compressed
    prs.save(dst)


def compress_docx(src: Path, dst: Path):
    doc = Document(src)
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            image_part = rel.target_part
            compressed = compress_image_bytes(image_part._blob)
            image_part._blob = compressed
    doc.save(dst)


def process_file(path: Path, test_only=False):
    if "_compressed" in path.stem:
        return None

    suffix = path.suffix.lower()
    if suffix not in (".pdf", ".pptx", ".docx"):
        return None

    dst = path.with_stem(path.stem + "_compressed")
    original_size = path.stat().st_size

    try:
        if suffix == ".pdf":
            compress_pdf(path, dst)
        elif suffix == ".pptx":
            compress_pptx(path, dst)
        elif suffix == ".docx":
            compress_docx(path, dst)
    except Exception as e:
        print(f"  ERRO em {path.name}: {e}")
        if dst.exists():
            dst.unlink()
        return None

    final_size = dst.stat().st_size
    reduction = (1 - final_size / original_size) * 100

    print(f"\n  Arquivo  : {path.name}")
    print(f"  Original : {fmt_size(original_size)}")
    print(f"  Final    : {fmt_size(final_size)}")
    print(f"  Redução  : {reduction:.1f}%")

    # Remove original only after confirmed success
    path.unlink()

    return original_size - final_size


def find_files(root: Path):
    files = []
    for dirpath, _, filenames in os.walk(root):
        for fname in filenames:
            p = Path(dirpath) / fname
            if p.suffix.lower() in (".pdf", ".pptx", ".docx") and "_compressed" not in p.stem:
                files.append(p)
    return sorted(files)


def main(test_only=False):
    if not BOOKS_DIR.exists():
        print(f"Pasta não encontrada: {BOOKS_DIR}")
        return

    files = find_files(BOOKS_DIR)

    if not files:
        print("Nenhum arquivo encontrado em public/books/")
        return

    if test_only:
        files = files[:1]
        print(f"=== MODO TESTE — processando apenas: {files[0].name} ===")
    else:
        print(f"=== Processando {len(files)} arquivo(s) ===")

    total_saved = 0
    processed = 0

    for f in files:
        saved = process_file(f)
        if saved is not None:
            total_saved += saved
            processed += 1

    print(f"\n{'='*45}")
    print(f"  Arquivos processados : {processed}")
    print(f"  Espaço economizado   : {fmt_size(total_saved)}")
    print(f"{'='*45}")


if __name__ == "__main__":
    test_mode = "--all" not in sys.argv
    main(test_only=test_mode)
